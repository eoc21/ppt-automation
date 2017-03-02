"""
Example usage:
entry('/Users/edwardcannon/Desktop/dataiku/res.csv','test1.pptx')

"""

import logging
import requests

from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

from pptx.util import Pt

import pandas as pd
import numpy as np

__author__ = 'edwardcannon'

logging.basicConfig(filename='ppt.log', level=logging.DEBUG)
LOGGER = logging.getLogger(__name__)

class InfluencerSlide(object):
    """
    Creates ppt slide for an
    influencer
    """
    def __init__(self, prs, data_row):
        self.layout = prs.slide_layouts[5]
        self.slide = prs.slides.add_slide(self.layout)
        self.shapes = self.slide.shapes
        self.data_row = data_row

    def iter_cells(self, table):
        """
        Iterate over cells in table
        :param table: Input table
        :return:
        """
        for row in table.rows:
            for cell in row.cells:
                yield cell

    def add_twitter_profile(self):
        """
        Add twitter profile information
        e.g. followers, handle etc
        :return:
        """
        rows = 2
        cols = 2
        left = Cm(0.81)
        top = Cm(2)
        width = Cm(11)
        height = Cm(2.73)
        table = self.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(5)
        table.columns[1].width = Cm(3)
        try:
            image_data = requests.get(self.data_row['twitter_profile_image']).content
            with open('image.jpg', 'wb') as test:
                test.write(image_data)
            self.shapes.add_picture('image.jpg', Cm(0.81), Cm(0.05))
        except requests.HTTPError:
            LOGGER.warn("failed to extract image!")
        table.cell(0, 0).text = 'Handle'
        table.cell(0, 1).text = str(self.data_row['twitter_name'])
        table.cell(1, 0).text = 'Followers'
        table.cell(1, 1).text = str(self.data_row['twitter_followers_count'])
        self.__resize_table_font(table, 12)

    def __resize_table_font(self, table, size):
        """
        Resize all cell fonts in a table
        :param table: Table to resize
        :param size: Font size
        :return:
        """
        for cell in self.iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(size)

    def add_account_details(self):
        """
        populates influencer
        account information
        :return:
        """
        rows = 4
        cols = 3
        left = Cm(0.81)
        top = Cm(6.06)
        width = Cm(10.17)
        height = Cm(3.73)
        table = self.shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Cm(3)
        table.columns[1].width = Cm(3)
        # write column headings
        table.cell(0, 1).text = 'Account Details'
        table.cell(0, 2).text = 'Channels'
        # write body cells
        table.cell(1, 0).text = 'Location'
        table.cell(2, 0).text = 'Time zone'
        table.cell(3, 0).text = 'Created @'
        table.cell(1, 1).text = str(self.data_row['twitter_location'])
        table.cell(2, 1).text = str(self.data_row['twitter_time_zone'])
        if str(self.data_row['twitter_created_at']):
            table.cell(3, 1).text = str(self.data_row['twitter_created_at'])
        table.cell(1, 2).text = 'Twitter'
        table.cell(2, 2).text = 'Youtube'
        table.cell(3, 2).text = 'Instagram'
        self.__resize_table_font(table, 12)

    def add_metrics_and_sentiment(self):
        """
        populates metrics table
        :return:
        """
        rows = 5
        cols = 2
        left = Cm(0.81)
        top = Cm(11.60)
        width = Cm(10.17)
        height = Cm(2.73)
        table = self.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(5)
        table.columns[1].width = Cm(2)
        table.cell(0, 0).text = 'Metric'
        table.cell(0, 1).text = 'Score'
        table.cell(1, 0).text = 'Average Reach'
        table.cell(2, 0).text = 'Positive Sentiment(%)'
        table.cell(3, 0).text = 'Negative Sentiment(%)'
        table.cell(4, 0).text = 'Average Impact'
        table.cell(1, 1).text = \
            str("{0:.2f}".format(self.data_row['twitter_audience_average_Reach']))
        table.cell(2, 1).text = \
            str("{0:.2f}".format(self.data_row['twitter_audience_positive_sentiment_percent']))
        table.cell(3, 1).text = \
            str("{0:.2f}".format(self.data_row['twitter_audience_negative_sentiment_percent']))
        table.cell(4, 1).text = \
            str("{0:.2f}".format(self.data_row['twitter_audience_average_Impact']))
        self.__resize_table_font(table, 12)

    def add_hindex(self):
        """
        Adds Twitter H-index
        :return:
        """
        top = Cm(15.99)
        left = Cm(0.81)
        width = Cm(3)
        height = Cm(1)
        tx_box = self.shapes.add_textbox(left, top, width, height)
        tframe = tx_box.text_frame.paragraphs[0]
        run = tframe.add_run()
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)
        font.bold = True
        font.italic = None
        font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
        run.text = "H-INDEX:"+str(self.data_row['hindex'])

    def add_gender_pie_chart(self):
        """
        Gender breakdown pie chart
        :return:
        """
        chart_data = ChartData()
        chart_data.categories = ['Male', 'Female']
        male_percent = self.data_row['twitter_audience_male_percent']
        female_percent = self.data_row['twitter_audience_female_percent']
        total = male_percent+female_percent
        m_percent = male_percent/total
        f_percent = female_percent/total
        chart_data.add_series('Series 1', (m_percent, f_percent))#sort out the % values
        x_coord, y_coord, cx_width, cy_height = Cm(10.46), Cm(0.44), Cm(7), Cm(6)
        chart = self.slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x_coord, y_coord, cx_width, cy_height, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'

        top = Cm(0.01)
        left = Cm(10.46)
        width = Cm(3)
        height = Cm(1)
        tx_box = self.shapes.add_textbox(left, top, width, height)
        t_frame = tx_box.text_frame
        t_frame.text = "Gender Breakdown"

    def add_audience_age_chart(self):
        """
        Creates audience age bar chart
        :return:
        """
        chart_data = ChartData()
        chart_data.categories = ['0-9', '10-17', '18-24',
                                 '25-34', '35-44', '45-54',
                                 '55-64', '65+']
        group_1 = self.data_row['0 to 9']
        group_2 = self.data_row['10 to 17']
        group_3 = self.data_row['18-24']
        group_4 = self.data_row['25-34']
        group_5 = self.data_row['35-44']
        group_6 = self.data_row['45-54']
        group_7 = self.data_row['55-64']
        group_8 = self.data_row['65+']
        chart_data.add_series('Series 1', (group_1, group_2,
                                           group_3, group_4,
                                           group_5, group_6,
                                           group_7, group_8))

        x_coord, y_coord, cx_width, cy_height = Cm(10.46), Cm(8.44), Cm(7), Cm(6)
        self.slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x_coord, y_coord, cx_width, cy_height, chart_data)

        tx_box = self.shapes.add_textbox(Cm(10.9), Cm(7.8), Cm(10.9), height=1)
        t_frame = tx_box.text_frame
        t_frame.text = "Age Distribution"


    def add_account_type_chart(self):
        """
        Gender breakdown pie chart
        :return:
        """
        chart_data = ChartData()
        chart_data.categories = ['Organisation', 'Individual']
        organisation_percent = self.data_row['twitter_audience_organisational']
        individual_percent = self.data_row['twitter_audience_individuals']
        total = organisation_percent+individual_percent
        i_percent = individual_percent/total
        o_percent = organisation_percent/total
        chart_data.add_series('Series 1', (i_percent, o_percent))
        x_coord, y_coord, cx_width, cy_height = Cm(18), Cm(0.44), Cm(7.9), Cm(6.9)
        chart = self.slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x_coord, y_coord, cx_width, cy_height, chart_data).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'

        top = Cm(0.01)
        left = Cm(18)
        width = Cm(3)
        height = Cm(1)
        t_box = self.shapes.add_textbox(left, top, width, height)
        t_frame = t_box.text_frame
        t_frame.text = "Account Type Distribution"


    def add_audience_interest_chart(self):
        """
        Adds audience interest bar chart
        :return:
        """
        chart_data = ChartData()
        interests = ['animals & pets', 'automotive',
                     'beauty/health & fitness',
                     'books', 'business',
                     'environment', 'family & parenting',
                     'fashion',	'fine arts',
                     'food & drinks', 'games',
                     'movies', 'music',
                     'photo & video', 'politics',
                     'science', 'shopping',
                     'sports', 'technology',
                     'travel', 'tv']
        interest_dict = {}
        for interest in interests:
            interest_dict[interest] = self.data_row[interest]
        import operator
        sorted_interests = sorted(interest_dict.items(), key=operator.itemgetter(1), reverse=True)
        top_5_interests = sorted_interests[0:5]
        interest_names = []
        values = []
        for name in top_5_interests:
            interest_names.append(name[0])
            if np.isnan(name[1]):
                values.append(0)
            else:
                values.append(name[1])

        chart_data.categories = interest_names
        chart_data.add_series('Series 1', (values))

        x_coord, y_coord, cx_width, cy_height = Cm(18), Cm(8.44), Cm(7), Cm(6)
        self.slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x_coord, y_coord, cx_width, cy_height, chart_data)

        tx_box = self.shapes.add_textbox(Cm(18), Cm(7.8), Cm(10.9), height=1)
        t_frame = tx_box.text_frame
        t_frame.text = "Interest Distribution"

    def add_instagram_metrics(self):
        """
        Added instagram stats
        :return:
        """
        rows = 3
        cols = 2
        left = Cm(10.46)
        top = Cm(14.60)
        width = Cm(8)
        height = Cm(3)
        table = self.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(5)
        table.columns[1].width = Cm(2)
        table.cell(0, 0).text = 'Instagram Metric'
        table.cell(0, 1).text = 'Score'
        table.cell(1, 0).text = 'Followers'
        table.cell(2, 0).text = 'Posts'
        table.cell(1, 1).text = str("{0:.2f}".format(self.data_row['instagram_followers']))
        table.cell(2, 1).text = str("{0:.2f}".format(self.data_row['instagram_posts']))
        self.__resize_table_font(table, 12)

    def __add_text(self, txt, x_coord, y_coord,
                   cx_width, cy_height):
        """
        add text box
        :param txt: Text to add
        :param x_coord: x-coordinate
        :param y_coord: y-coordinate
        :param cx_width: x-width
        :param cy_height: y-height
        :return:
        """
        tx_box = self.shapes.add_textbox(Cm(x_coord), Cm(y_coord), Cm(cx_width), height=cy_height)
        t_frame = tx_box.text_frame
        t_frame.text = txt

    def add_youtube_metrics(self):
        """
        Add youtube stats
        :return:
        """
        rows = 5
        cols = 2
        left = Cm(18)
        top = Cm(14.60)
        width = Cm(8)
        height = Cm(3)
        if self.data_row['yt_channel_video_count'] != 0:
            rows = 4
        table = self.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(5)
        table.columns[1].width = Cm(2)
        table.cell(0, 0).text = 'Youtube Metric'
        table.cell(0, 1).text = 'Score'
        if self.data_row['yt_channel_video_count'] == 0:
            table.cell(1, 0).text = 'Average Comments'
            table.cell(2, 0).text = 'Average Dislikes'
            table.cell(3, 0).text = 'Average Likes'
            table.cell(4, 0).text = 'Average Views'
            table.cell(1, 1).text = str("{0:.2f}".format(self.data_row['youtube_avg_comments']))
            table.cell(2, 1).text = str("{0:.2f}".format(self.data_row['youtube_avg_dislikes']))
            table.cell(3, 1).text = str("{0:.2f}".format(self.data_row['youtube_avg_likes']))
            table.cell(4, 1).text = str("{0:.2f}".format(self.data_row['youtube_avg_views']))
        else:
            table.cell(1, 0).text = 'Channel Comments'
            table.cell(2, 0).text = 'Channel Videos'
            table.cell(3, 0).text = 'Channel Views'
            table.cell(1, 1).text = str("{0:.2f}".format(self.data_row['yt_channel_comment_count']))
            table.cell(2, 1).text = str("{0:.2f}".format(self.data_row['yt_channel_video_count']))
            table.cell(3, 1).text = str("{0:.2f}".format(self.data_row['yt_channel_view_count']))
        self.__resize_table_font(table, 12)

    def add_hindex_footer(self):
        """
        Adds H-index footer
        :return:
        """
        text = 'H-Index is a measure of productivity (tweets) '\
        '& engagement (retweets). An influencer has an '\
        'index of H when they have posted H tweets that '\
        'has received at least H retweets'
        self.__add_text(text, 0.05, 18, 18, 1)

def entry(input_file, output):
    """
    ppt generate entry point
    :param input_file: input csv with influencer metrics
    :param output: output pptx
    :return:
    """
    input_df = pd.read_csv(input_file)
    prs = Presentation()
    for index, row in input_df.iterrows():
        influencer = InfluencerSlide(prs, row)
        influencer.add_account_details()
        influencer.add_metrics_and_sentiment()
        influencer.add_hindex()
        influencer.add_gender_pie_chart()
        influencer.add_account_type_chart()
        influencer.add_audience_age_chart()
        influencer.add_audience_interest_chart()
        influencer.add_instagram_metrics()
        influencer.add_youtube_metrics()
        influencer.add_twitter_profile()
    prs.save(output)

if __name__ == '__main__':
    pass
